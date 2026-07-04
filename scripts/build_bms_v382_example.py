"""
v3.8.2 视觉填满 + 文字框紧凑
==========================
老大的两个新规则：
1. 图片文字框虽然铺满，但视觉上文字整体在左边 — 必须把"云 BMS"段挪到右边
2. 文字框大小要紧凑 — 框内没有多余空间，然后对框进行布局

实现：
- 传统 BMS 在左（框 3.55 寸宽）
- 云 BMS 在右（框 3.55 寸宽）
- 每个文字框按内容紧凑定高（用 calc_h 算最小高度）
- 标题带在 2 个框上方
- 文字框内 margin 最小化（0.02"），无多余空间
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from PIL import Image

PAGE_W, PAGE_H = 13.33, 7.50

COLOR_MAIN = RGBColor(0x0E, 0x3A, 0x6B)
COLOR_ACCENT = RGBColor(0xC0, 0x39, 0x2B)
COLOR_TEXT = RGBColor(0x2C, 0x3E, 0x50)
COLOR_BLACK = RGBColor(0x00, 0x00, 0x00)
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_BG = RGBColor(0xF5, 0xF7, 0xFA)
COLOR_BORDER = RGBColor(0xD5, 0xD9, 0xE0)


def calc_h(text, font_pt, line_spacing=1.20, padding=0.04):
    """v3.7.5 真实行高公式 — 文字框紧凑定高"""
    lines = sum(1 for p in text.split('\n') if p.strip())
    return lines * font_pt * line_spacing / 72.0 + padding * 2


def add_text_compact(slide, x, y, text, font_pt=10, bold=False, color=COLOR_TEXT,
                     align=PP_ALIGN.LEFT, bg_color=None, line_spacing=1.20,
                     min_padding=0.04, fixed_w=None):
    """紧凑文字框 — 框高按内容定，无多余空间"""
    h = calc_h(text, font_pt, line_spacing, min_padding)
    # 估算框宽 — 按最长行 0.18 寸/字符
    max_line = max(len(l) for l in text.split('\n') if l.strip())
    w = fixed_w if fixed_w else min(max_line * 0.18 + 0.1, 4.0)
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.02); tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(min_padding); tf.margin_bottom = Inches(min_padding)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    if bg_color is not None:
        tb.fill.solid(); tb.fill.fore_color.rgb = bg_color
        tb.line.fill.background()
    lines = text.split('\n')
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.line_spacing = line_spacing
        run = p.add_run(); run.text = line
        run.font.size = Pt(font_pt); run.font.bold = bold
        run.font.color.rgb = color; run.font.name = 'Microsoft YaHei'
    return tb, w, h


def add_picture_iso(slide, img_path, x, y, max_w, max_h):
    img = Image.open(img_path)
    iw, ih = img.size
    img_ratio = iw / ih
    calc_h = max_h
    calc_w = calc_h * img_ratio
    if calc_w > max_w:
        calc_w = max_w
        calc_h = calc_w / img_ratio
    actual_x = x + (max_w - calc_w) / 2
    actual_y = y + (max_h - calc_h) / 2
    return slide.shapes.add_picture(img_path, Inches(actual_x), Inches(actual_y),
                                     Inches(calc_w), Inches(calc_h))


def add_band(slide, x, y, w, h, text, color, font_pt=13, font_color=COLOR_WHITE):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid(); shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.margin_left = Inches(0.1); tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.04); tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    run = p.add_run(); run.text = text
    run.font.size = Pt(font_pt); run.font.bold = True
    run.font.color.rgb = font_color; run.font.name = 'Microsoft YaHei'


def add_figure_caption(slide, x, y, w, text):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(0.25))
    tf = tb.text_frame
    tf.margin_left = Inches(0.02); tf.margin_right = Inches(0.02)
    tf.margin_top = Inches(0.02); tf.margin_bottom = Inches(0.02)
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    run = p.add_run(); run.text = text
    run.font.size = Pt(10); run.font.bold = True
    run.font.color.rgb = COLOR_BLACK; run.font.name = 'Microsoft YaHei'


# ============= 主程序 =============
prs = Presentation()
prs.slide_width = Inches(PAGE_W); prs.slide_height = Inches(PAGE_H)
blank = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank)

bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(PAGE_W), Inches(PAGE_H))
bg.fill.solid(); bg.fill.fore_color.rgb = COLOR_WHITE
bg.line.fill.background()

# 1. 主标题 — 用紧凑 calc_h
title_text = "0  概念与对比\n云 BMS vs 传统 BMS"
TITLE_H = calc_h(title_text, 22, 1.10) + 0.10
add_text_compact(slide, 0.20, 0.05, title_text,
                 font_pt=22, bold=True, color=COLOR_MAIN, line_spacing=1.10, min_padding=0.05)

# 副标题带（红色横条）
sub_y = 0.05 + TITLE_H + 0.05
SUB_H = 0.40
sub_band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.20), Inches(sub_y),
                                    Inches(13.00), Inches(SUB_H))
sub_band.fill.solid(); sub_band.fill.fore_color.rgb = COLOR_ACCENT
sub_band.line.fill.background()
tf = sub_band.text_frame
tf.margin_left = Inches(0.15); tf.margin_top = Inches(0.05)
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
run = p.add_run()
run.text = "端-边-云协同架构重塑电池管理系统  ·  算力 1000×  ·  预警 16 天  ·  软件市场 CAGR 20.6%"
run.font.size = Pt(11); run.font.bold = True
run.font.color.rgb = COLOR_WHITE; run.font.name = 'Microsoft YaHei'

# ===== v3.8.2 关键：上排左 = 左右分栏（传统左 + 云右）=====
UPPER_Y = sub_y + SUB_H + 0.20
LEFT_X = 0.20
LEFT_W = 7.50  # 总宽
HALF_W = 3.65  # 每栏宽

# 标题带（占整个左区域）
add_band(slide, LEFT_X, UPPER_Y, LEFT_W, 0.40,
         "一、定义 · 端-边-云 3 层架构", COLOR_MAIN, font_pt=14)

# ===== 关键改动 1：左栏 = 传统 BMS（紧凑框）=====
LEFT_HALF_X = LEFT_X + 0.05
trad_title = "▶ 传统 BMS（电池端 MCU）"
trad_body = "● 单 MCU + AFE · 固定阈值告警\n● 算力 ~10 MIPS, 存储 KB 级\n● 无 AI 模型，无云端协同"

# 标题（紧凑）
tb, tw, th = add_text_compact(slide, LEFT_HALF_X, UPPER_Y + 0.50, trad_title,
                              font_pt=12, bold=True, color=COLOR_MAIN, line_spacing=1.20)
# 内容（紧凑）— 框紧贴标题下方
body_y = UPPER_Y + 0.50 + th + 0.03
tb2, tw2, th2 = add_text_compact(slide, LEFT_HALF_X, body_y, trad_body,
                                 font_pt=10, color=COLOR_TEXT, line_spacing=1.20)

# ===== 关键改动 2：右栏 = 云 BMS（紧凑框）=====
RIGHT_HALF_X = LEFT_X + HALF_W + 0.10
cloud_title = "▶ 云 BMS（端 + 边 + 云）"
cloud_body = "● 云端算力 + AI 大模型 · 实时状态估计 + 故障预测\n● 算力 10,000+ MIPS, 存储 TB 级\n● 跨车队数据闭环，OTA 持续迭代"

# 标题（紧凑）
tb3, tw3, th3 = add_text_compact(slide, RIGHT_HALF_X, UPPER_Y + 0.50, cloud_title,
                                 font_pt=12, bold=True, color=COLOR_ACCENT, line_spacing=1.20)
# 内容（紧凑）
body_y2 = UPPER_Y + 0.50 + th3 + 0.03
tb4, tw4, th4 = add_text_compact(slide, RIGHT_HALF_X, body_y2, cloud_body,
                                 font_pt=10, color=COLOR_TEXT, line_spacing=1.20)

# 上排总高 = 标题带 + 标题 + 内容 + 间距
UPPER_H = 0.40 + 0.10 + max(th + th2, th3 + th4) + 0.30

# 2. 上排右：核心数据（保持原样）
RIGHT_X = 7.80
RIGHT_W = 5.40
add_text_compact(slide, RIGHT_X, UPPER_Y, "★ 核心数据",
                 font_pt=14, bold=True, color=COLOR_ACCENT, line_spacing=1.20, min_padding=0.04)
data_items = [
    ("算力提升", "1000×", "10 → 10,000+ MIPS"),
    ("存储提升", "10⁶×", "KB → TB 级"),
    ("预警提前", "16 天", "热失控预测"),
    ("市场 CAGR", "20.6%", "2026-2032"),
]
item_h = 0.42
for i, (label, big, sub) in enumerate(data_items):
    iy = UPPER_Y + 0.35 + i * item_h
    add_text_compact(slide, RIGHT_X, iy, label,
                     font_pt=10, bold=True, color=COLOR_MAIN, line_spacing=1.10, fixed_w=1.5)
    add_text_compact(slide, RIGHT_X + 1.5, iy, big,
                     font_pt=14, bold=True, color=COLOR_ACCENT, line_spacing=1.10,
                     align=PP_ALIGN.CENTER, fixed_w=1.5)
    add_text_compact(slide, RIGHT_X + 3.0, iy, sub,
                     font_pt=9, color=COLOR_TEXT, line_spacing=1.10, fixed_w=2.4)

# 3. 下排左：6 项差异
LOWER_Y = UPPER_Y + UPPER_H + 0.15
LOWER_H = 2.40
add_band(slide, LEFT_X, LOWER_Y, LEFT_W, 0.40,
         "二、6 项本质差异", COLOR_MAIN, font_pt=13)

cards_y = LOWER_Y + 0.45
card_h = 0.62
card_w = 3.65
card_gap_x = 0.10
card_gap_y = 0.05

cards = [
    ("① 算力", "10 MIPS → 10,000+ MIPS", "(+1000×)"),
    ("② 存储", "KB 级 → TB 级", "(+10⁶×)"),
    ("③ 时延", "ms 级本地 → s 级云端协同", ""),
    ("④ 模型更新", "出厂固化 → OTA 持续迭代", ""),
    ("⑤ 业务闭环", "单车数据 → 跨车队", ""),
    ("⑥ 商业模式", "卖硬件 → 软件订阅", ""),
]

for idx, (label, val, ext) in enumerate(cards):
    row = idx // 2
    col = idx % 2
    cx = LEFT_X + col * (card_w + card_gap_x)
    cy = cards_y + row * (card_h + card_gap_y)
    card_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(cx), Inches(cy),
                                      Inches(card_w), Inches(card_h))
    card_bg.fill.solid(); card_bg.fill.fore_color.rgb = COLOR_WHITE
    card_bg.line.color.rgb = COLOR_BORDER
    card_bg.line.width = Pt(0.75)
    add_text_compact(slide, cx + 0.05, cy + 0.04, label,
                     font_pt=10, bold=True, color=COLOR_MAIN, line_spacing=1.10, fixed_w=card_w-0.10)
    add_text_compact(slide, cx + 0.05, cy + 0.26, val,
                     font_pt=9, color=COLOR_TEXT, line_spacing=1.10, fixed_w=card_w-0.10)
    if ext:
        add_text_compact(slide, cx + 0.05, cy + 0.46, ext,
                         font_pt=8, bold=True, color=COLOR_ACCENT, line_spacing=1.10,
                         fixed_w=card_w-0.10)

# 4. 下排右：图
LOWER_RIGHT_X = 7.80
LOWER_RIGHT_W = 5.40

# 雷达图
fig1_y = LOWER_Y
fig1_w = 4.0
fig1_h = fig1_w * (10/16)
add_picture_iso(slide, 'artifacts/charts/01_雷达图_6维对比.png',
                LOWER_RIGHT_X, fig1_y, fig1_w, fig1_h)
add_figure_caption(slide, LOWER_RIGHT_X, fig1_y + 2.55, LOWER_RIGHT_W,
                   "图 1  6 维评分对比 — 云 BMS 全面领先")

# 架构图
fig2_y = LOWER_Y + 2.90
fig2_w = 4.5
fig2_h = fig2_w * (6/10)
add_picture_iso(slide, 'artifacts/charts/06_架构图_端边云.png',
                LOWER_RIGHT_X, fig2_y, fig2_w, fig2_h)
add_figure_caption(slide, LOWER_RIGHT_X, fig2_y + 2.75, LOWER_RIGHT_W,
                   "图 2  端-边-云 4 能力 × 3 层架构")

# 5. 底部引言
BOTTOM_Y = LOWER_Y + LOWER_H + 0.10
BOTTOM_H = 0.40
add_text_compact(slide, 0.20, BOTTOM_Y,
                 "云 BMS 本质 = 云端算力 + AI 模型 + 跨资产数据 + 数字孪生  ·  出处: RSC 2025 / 调研报告",
                 font_pt=10, color=COLOR_TEXT, line_spacing=1.20,
                 bg_color=COLOR_BG, fixed_w=13.00)

out = 'artifacts/BMS_vs_云BMS_对比页_v3.8.2_视觉填满+紧凑框.pptx'
prs.save(out)
print(f'OK Saved: {out}')
print(f'  上排左总高: {UPPER_H:.2f} (含标题带+左右双栏)')
print(f'  左栏总高: {th + th2:.2f}, 右栏总高: {th3 + th4:.2f}')
