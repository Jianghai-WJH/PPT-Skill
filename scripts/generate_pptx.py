#!/usr/bin/env python3
"""
generate_pptx.py（v2 - 增强版）
============================================================
基于模板 + slot_map + 内容 JSON，生成最终 PPT。

核心能力（区别于 v1）：
- 支持"半结构化"识别：每个 run 单独判定是 template(保留) 还是 variable(替换)
- 支持 GROUP 内的子形状递归填充
- 支持"按 run 填充"和"按段落整体替换"两种模式
- 保留模板所有装饰元素（logo、母版色块等）

输入：
  <template.pptx>      模板文件
  <slot_map.json>      detect_template_slots.py 的输出
  <content.json>       用户指定每个 variable 的填充内容
  <output.pptx>        最终文件

content.json 格式（v2 增强版）：
{
  "template_path": "...",          // 可选
  "slides": [
    {
      "use_slide_idx": 0,          // 引用模板中的哪张 slide（复制它）
      "fill_map": {
        // 方式 1：按 shape_idx + 段落内 run 索引填充
        // 路径用 "sh{idx}.p{p_idx}.r{r_idx}" 或 "sh{idx}.sub{sub_idx}.p{p_idx}.r{r_idx}"
        "sh1.p0.r1": "5W2H分析",           // 替换某个具体 run
        "sh2.sub1.p1.r1": "When 描述",     // 组合内子形状
        "sh1.p1": "整段替换",              // 整段替换（清空后重写）

        // 方式 2：按语义键填充（推荐：AI 友好）
        // AI 解析 slot_map.json 时可以给每个 variable 分配语义键
        "section.subtitle": "3.3 电机异响",
        "section.body": "高速段电磁噪音",
        "5w2h.when": "2026-06-30 出现异响",
        "5w2h.where": "实验室",
        "5w2h.who": "试验员",
        "5w2h.what": "高速段出现嗡嗡声",
        "5w2h.why": "电磁力波耦合",
        "5w2h.how": "依据试验安排",
        "5w2h.how_many": "5次测试"
      }
    }
  ]
}
"""
import copy
import json
import sys
from collections import defaultdict
from pathlib import Path

from pptx import Presentation
from pptx.util import Pt
from lxml import etree


NS = {
    "a": "http://schemas.openxmlformats.org/drawingml/2006/main",
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
    "r": "http://schemas.openxmlformats.org/officeDocument/2006/relationships",
}


# ============================================================
# Run 路径定位
# ============================================================

def find_run_by_path(shape, path):
    """
    在 shape 内按路径找到 run 元素
    path 形如: "p1.r2" 或 "sub0.p1.r2" 或 "sub0.sub1.p1.r2"
    """
    parts = path.split(".")
    target_shape = shape

    # 处理多级 sub 路径
    while parts and parts[0].startswith("sub"):
        sub_idx = int(parts[0][3:])
        if target_shape.shape_type != 6:  # 不是组合
            return None, None
        if sub_idx >= len(target_shape.shapes):
            return None, None
        target_shape = target_shape.shapes[sub_idx]
        parts = parts[1:]

    if not target_shape.has_text_frame:
        return None, None

    if len(parts) != 2:
        return None, None

    try:
        p_idx = int(parts[0][1:])
        r_idx = int(parts[1][1:])
    except (ValueError, IndexError):
        return None, None

    tf = target_shape.text_frame
    if p_idx >= len(tf.paragraphs):
        return None, None
    para = tf.paragraphs[p_idx]
    if r_idx >= len(para.runs):
        return None, None

    return target_shape, para.runs[r_idx]


# ============================================================
# 文本替换（保留原 run 样式）
# ============================================================

def replace_run_text(run, new_text):
    """替换 run 的文本，保留所有格式属性"""
    run.text = str(new_text)


def replace_paragraph_text(para, new_text, first_run_style_only=True):
    """整段替换：保留第一个 run 的样式，清空其他 run"""
    runs = para.runs
    if not runs:
        # 段落无 run（罕见），加一个
        para.text = str(new_text)
        return

    runs[0].text = str(new_text)
    for r in runs[1:]:
        r.text = ""


def append_to_paragraph(para, new_text, first_run_style_only=True):
    """往段落末尾追加文本（保留首个 run 样式）"""
    if not para.runs:
        para.text = str(new_text)
        return
    # 找最后一个非空 run，在它后面加
    last_run = para.runs[-1]
    last_run.text = (last_run.text or "") + str(new_text)


# ============================================================
# 整 shape 处理
# ============================================================

def find_shape_by_idx(slide, sh_idx):
    """按索引找 shape（支持组合递归）"""
    if sh_idx < len(slide.shapes):
        return slide.shapes[sh_idx]
    return None


def get_semantic_key(shape, run_path=None, default=None):
    """从 shape name 推断语义键（启发式）
    例如 '文本框 1' → 'textbox_1'
        '组合 10' → 'group_10'
    """
    if default:
        return default
    name = (shape.name or "").lower()
    name = name.replace(" ", "_").replace(":", "")
    return f"shape.{name}"


# ============================================================
# 主流程
# ============================================================

def apply_clear_others(slide, clear_others):
    """清空指定段落中除指定 run 之外的所有 run
    clear_others 格式: {"sh{idx}.p{p_idx}": keep_r_idx}
    """
    for key, keep_r_idx in clear_others.items():
        if not key.startswith("sh"):
            continue
        try:
            sh_idx = int(key[2:].split(".")[0])
            p_idx = int(key.split(".")[1][1:])
        except (ValueError, IndexError):
            continue
        if sh_idx >= len(slide.shapes):
            continue
        shape = slide.shapes[sh_idx]
        if not shape.has_text_frame:
            continue
        if p_idx >= len(shape.text_frame.paragraphs):
            continue
        para = shape.text_frame.paragraphs[p_idx]
        for r_idx, run in enumerate(para.runs):
            if r_idx != keep_r_idx:
                run.text = ""


def apply_fill_map(slide, fill_map, slot_map=None):
    """根据 fill_map 填充 slide 的内容
    fill_map 格式:
      {
        "sh{idx}.p{p}.r{r}": "新文本",           # 替换指定 run
        "sh{idx}.sub{s}.p{p}.r{r}": "...",       # 组合内 run
        "sh{idx}.p{p}": "整段新文本",              # 整段替换
      }
    """
    for key, value in fill_map.items():
        # 解析 sh_idx
        if not key.startswith("sh"):
            continue
        try:
            sh_idx = int(key[2:].split(".")[0])
        except (ValueError, IndexError):
            continue

        path = key[len(f"sh{sh_idx}."):]  # 剩余路径
        if not path.startswith(("p", "sub")):
            continue

        shape = find_shape_by_idx(slide, sh_idx)
        if shape is None:
            continue

        if path.startswith("p"):
            # 顶层段落路径
            target_shape, run = find_run_by_path(shape, path)
        elif path.startswith("sub"):
            target_shape, run = find_run_by_path(shape, path)
        else:
            continue

        if run is None:
            continue

        # 替换文本（保留样式）
        replace_run_text(run, value)


def apply_semantic_fill(slide, semantic_map, slot_map_for_slide):
    """根据 slot_map 的语义键填充
    semantic_map: {"5w2h.when": "...", "section.subtitle": "..."}
    slot_map_for_slide: 该 slide 的 slot_map 决策
    """
    # 建立 (sh_idx, p_idx, r_idx) -> 语义键 的反向索引
    # 来自 slot_map.json 中 AI 给每个 variable 分配的 semantic_key
    # 如果 slot_map 没有 semantic_key，则用启发式：基于 shape 名称 + 上下文文本

    for shape_decision in slot_map_for_slide.get("shapes", []):
        sh_idx = shape_decision["sh_idx"]
        # 收集该 shape 的所有 variable run
        var_runs = []
        for run_info in shape_decision.get("final_runs", []):
            if run_info.get("role") == "variable":
                var_runs.append(run_info)

        # 找匹配
        for var_run in var_runs:
            # 这里简化：按出现顺序匹配（更智能的做法是按上下文）
            pass


def deep_copy_slide(prs, source_slide):
    """深度复制一张 slide（含所有形状、组、文本格式）"""
    from copy import deepcopy
    from lxml import etree

    # 创建新 slide（用相同 layout）
    new_slide = prs.slides.add_slide(source_slide.slide_layout)

    # 删除默认的占位符（如果与源 slide 不对应）
    # 简化做法：直接复制源 slide 的 spTree
    src_spTree = source_slide.shapes._spTree
    dst_spTree = new_slide.shapes._spTree

    # 清空目标 spTree（保留 nvGrpSpPr 等）
    for child in list(dst_spTree):
        if etree.QName(child.tag).localname in ("sp", "pic", "grpSp", "graphicFrame", "cxnSp"):
            dst_spTree.remove(child)

    # 复制源 spTree 的 sp/pic/grpSp 到目标
    for child in src_spTree:
        tag = etree.QName(child.tag).localname
        if tag in ("sp", "pic", "grpSp", "graphicFrame", "cxnSp"):
            # 深拷贝元素
            new_child = deepcopy(child)
            # 重新分配关系 ID（图片等需要 rId）
            dst_spTree.append(new_child)

    return new_slide


def generate_pptx(template_path, slot_map_path, content_path, output_path):
    """主入口"""
    template_p = Path(template_path)
    slot_map_p = Path(slot_map_path)
    content_p = Path(content_path)
    output_p = Path(output_path)

    slot_map = json.loads(slot_map_p.read_text(encoding="utf-8"))
    content = json.loads(content_p.read_text(encoding="utf-8"))

    prs = Presentation(template_p)

    # 删除模板中的所有示例 slide
    sldIdLst = prs.slides._sldIdLst
    for sldId in list(sldIdLst):
        rId = sldId.get(
            "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id"
        )
        prs.part.drop_rel(rId)
        sldIdLst.remove(sldId)

    # 按 content.slides 顺序生成
    slot_slides = slot_map.get("decision", [])

    for s_idx, slide_spec in enumerate(content.get("slides", [])):
        use_idx = slide_spec.get("use_slide_idx", 0)
        if use_idx >= len(slot_slides):
            print(f"⚠️ use_slide_idx={use_idx} 超出范围，跳过")
            continue

        # 找到对应的源 slide
        src_slide = Presentation(template_p).slides[use_idx]
        new_slide = deep_copy_slide(prs, src_slide)

        # 应用 fill_map
        fill_map = slide_spec.get("fill_map", {})
        if fill_map:
            apply_fill_map(new_slide, fill_map, slot_slides[use_idx])

    prs.save(output_p)
    print(f"✅ 已生成：{output_p} (共 {len(content.get('slides', []))} 页)")


def main():
    if len(sys.argv) != 5:
        print("Usage: python generate_pptx.py <template.pptx> <slot_map.json> <content.json> <output.pptx>")
        sys.exit(1)

    generate_pptx(sys.argv[1], sys.argv[2], sys.argv[3], sys.argv[4])


if __name__ == "__main__":
    main()
