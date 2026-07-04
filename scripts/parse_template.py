#!/usr/bin/env python3
"""
parse_template.py
解析用户提供的 PPT 模板，提取版式、字体、字号、配色、装饰元素等元数据。
输出为 JSON 文件，供后续生成步骤使用。

Usage:
    python parse_template.py <template.pptx> <output.json>
"""
import json
import sys
from pathlib import Path
from collections import Counter

from pptx import Presentation
from pptx.util import Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor


def emu_to_pt(emu):
    """EMU -> pt（仅用于尺寸参考展示）"""
    if emu is None:
        return None
    return round(emu / 12700, 2)


def extract_color(rgb_or_theme):
    """统一处理颜色对象 -> hex 字符串"""
    if rgb_or_theme is None:
        return None
    try:
        return "#" + str(rgb_or_theme).upper()
    except Exception:
        return None


def font_summary(font):
    """从 Font 对象提取摘要信息"""
    if font is None:
        return None
    return {
        "name": font.name,
        "size_pt": float(font.size.pt) if font.size else None,
        "bold": font.bold,
        "italic": font.italic,
    }


def extract_placeholder(placeholder):
    """提取占位符信息"""
    ph_format = placeholder.placeholder_format
    info = {
        "idx": ph_format.idx,
        "type": str(ph_format.type) if ph_format.type else None,
        "name": placeholder.name,
        "position": {
            "left_pt": emu_to_pt(placeholder.left),
            "top_pt": emu_to_pt(placeholder.top),
            "width_pt": emu_to_pt(placeholder.width),
            "height_pt": emu_to_pt(placeholder.height),
        },
    }
    # 尝试提取示例文本与字体
    if placeholder.has_text_frame:
        tf = placeholder.text_frame
        sample_text = tf.text.strip()
        info["sample_text"] = sample_text if sample_text else None
        if tf.paragraphs and tf.paragraphs[0].runs:
            run = tf.paragraphs[0].runs[0]
            info["font"] = font_summary(run.font)
    return info


def extract_shape(shape):
    """提取形状信息（简化版，用于识别装饰元素/logo）"""
    info = {
        "name": shape.name,
        "type": str(shape.shape_type) if shape.shape_type else None,
        "position": {
            "left_pt": emu_to_pt(shape.left),
            "top_pt": emu_to_pt(shape.top),
            "width_pt": emu_to_pt(shape.width),
            "height_pt": emu_to_pt(shape.height),
        },
    }
    # 是否包含文字
    if shape.has_text_frame:
        text = shape.text_frame.text.strip()
        if text:
            info["text"] = text
    # 是否为图片（识别为 logo / 装饰图）
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        info["is_picture"] = True
    # 形状填充色
    try:
        if hasattr(shape, "fill") and shape.fill.type is not None:
            if shape.fill.type == 1:  # solid
                info["fill_color"] = extract_color(shape.fill.fore_color.rgb)
    except Exception:
        pass
    return info


def classify_layout(layout_name, placeholders):
    """基于名称+占位符模式推断版式类型"""
    name_lower = (layout_name or "").lower()
    has_title = any(p["type"] and "TITLE" in str(p["type"]) for p in placeholders)
    has_body = any(p["type"] and "BODY" in str(p["type"]) for p in placeholders)
    ph_count = len(placeholders)

    # 启发式分类（顺序敏感：cover 在前以避免 "title slide" 命中 toc）
    if "cover" in name_lower:
        return "cover"
    if "目录" in name_lower or "agenda" in name_lower or "table of contents" in name_lower or name_lower.strip() == "toc":
        return "toc"
    if "summary" in name_lower or "conclusion" in name_lower or "谢谢" in name_lower or "thank" in name_lower:
        return "summary"
    if "title" in name_lower and ph_count <= 1:
        return "cover"
    if "section" in name_lower:
        return "section_header"
    if has_title and has_body:
        return "content"
    if has_title and not has_body:
        return "section_header"
    return "other"


def parse_template(template_path):
    """主函数：解析 PPT 模板"""
    prs = Presentation(template_path)

    # 1. 全局尺寸
    global_info = {
        "slide_width_pt": emu_to_pt(prs.slide_width),
        "slide_height_pt": emu_to_pt(prs.slide_height),
    }

    # 2. 主题色与字体方案（从母版中提取）
    theme = {}
    try:
        if prs.slide_masters:
            master = prs.slide_masters[0]
            # 字体方案
            try:
                theme["major_font"] = (
                    master.element.xpath("//a:majorFont/a:latin")
                    and master.element.xpath("//a:majorFont/a:latin")[0]
                    .get("typeface")
                )
            except Exception:
                pass
            try:
                theme["minor_font"] = (
                    master.element.xpath("//a:minorFont/a:latin")
                    and master.element.xpath("//a:minorFont/a:latin")[0]
                    .get("typeface")
                )
            except Exception:
                pass
            # 颜色方案
            try:
                clr_scheme = master.element.xpath("//a:clrMap")
                if clr_scheme:
                    theme["color_map"] = dict(
                        (k, v) for k, v in clr_scheme[0].attrib.items()
                    )
            except Exception:
                pass
    except Exception:
        pass

    # 3. 母版级版式
    layouts_info = []
    for layout in prs.slide_layouts:
        placeholders = []
        for ph in layout.placeholders:
            placeholders.append(extract_placeholder(ph))
        # 版式中的非占位符形状（装饰、logo、背景）
        shapes = []
        for shape in layout.shapes:
            if shape.is_placeholder:
                continue
            shapes.append(extract_shape(shape))
        layouts_info.append(
            {
                "name": layout.name,
                "type": classify_layout(layout.name, placeholders),
                "placeholders": placeholders,
                "decorative_shapes": shapes,
            }
        )

    # 4. 模板中已有的示例页（如果有）
    sample_slides = []
    for idx, slide in enumerate(prs.slides):
        if idx >= 20:  # 防止异常模板过多
            break
        slide_info = {
            "idx": idx,
            "layout_name": slide.slide_layout.name,
            "shapes": [extract_shape(s) for s in slide.shapes],
        }
        sample_slides.append(slide_info)

    # 5. 汇总：统计字体、字号、颜色
    all_fonts = Counter()
    all_sizes = Counter()
    all_colors = Counter()
    for layout in layouts_info:
        for ph in layout["placeholders"]:
            font = ph.get("font")
            if font:
                if font["name"]:
                    all_fonts[font["name"]] += 1
                if font["size_pt"]:
                    all_sizes[int(font["size_pt"])] += 1
        for shape in layout["decorative_shapes"]:
            if shape.get("fill_color"):
                all_colors[shape["fill_color"]] += 1

    stats = {
        "fonts_top5": all_fonts.most_common(5),
        "sizes_top10": all_sizes.most_common(10),
        "colors_top10": [(c, n) for c, n in all_colors.most_common(10)],
    }

    return {
        "source": str(Path(template_path).resolve()),
        "global": global_info,
        "theme": theme,
        "layouts": layouts_info,
        "sample_slides": sample_slides,
        "stats": stats,
    }


def main():
    if len(sys.argv) != 3:
        print("Usage: python parse_template.py <template.pptx> <output.json>")
        sys.exit(1)
    template_path = sys.argv[1]
    output_path = sys.argv[2]
    result = parse_template(template_path)
    with open(output_path, "w", encoding="utf-8") as f:
        json.dump(result, f, ensure_ascii=False, indent=2)
    print(f"✅ 解析完成：{output_path}")
    print(f"   版式数：{len(result['layouts'])}")
    print(f"   主题字体：{result['theme'].get('major_font')} / {result['theme'].get('minor_font')}")
    print(f"   Top 字号：{[s[0] for s in result['stats']['sizes_top10'][:5]]}")


if __name__ == "__main__":
    main()
