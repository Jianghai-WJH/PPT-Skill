#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
云枢·v3.7.1 修复版
- 图 1 给足 1:1 空间（限制 max_h=1.85 max_w=2.20, 取 min）
- 修 5 处文字堆叠：
  * 云 BMS 末行与"二、6 项"标题带重叠 -> 缩短云 BMS 区域高
  * "★核心数据"textbox 与数据点 textbox 重叠 -> 合并
  * 卡片 "+1000×" / "+10⁶×" 与"算力"行重叠 -> 调整 card 高度
"""

import os
import shutil
from PIL import Image
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

SRC = "/home/gem/.aily/workdir/web_p2p_b8309bbd/attachments/AI云BMS-五看 - 副本.pptx"
OUT = "/home/gem/.aily/workdir/web_p2p_b8309bbd/BMS_vs_云BMS_对比页_v3.7.1_空间自适应修复.pptx"
CHART_RADAR = "/home/gem/.aily/workdir/web_p2p_b8309bbd/charts/01_雷达图_6维对比.png"
CHART_ARCH = "/home/gem/.aily/workdir/web_p2p_b8309bbd/charts/06_架构图_端边云.png"

PAGE_W = 13.333
PAGE_H = 7.5

COLOR_MAIN = RGBColor(0x1F, 0x4E, 0x79)
COLOR_RED = RGBColor(0xC0, 0x00, 0x00)
COLOR_GRAY = RGBColor(0x40, 0x40, 0x40)
COLOR_LIGHT_GRAY = RGBColor(0x8C, 0x8C, 0x8C)
COLOR_HIGHLIGHT = RGBColor(0xE9, 0x71, 0x32)
COLOR_BG_LIGHT = RGBColor(0xF2, 0xF2, 0xF2)
COLOR_BG_RED = RGBColor(0xFD, 0xEC, 0xEC)
COLOR_BG_BLUE = RGBColor(0xE8, 0xEF, 0xF7)
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def set_run(run, text, size=10, bold=False, color=COLOR_GRAY, italic=False):
    run.text = text
    run.font.name = 'Microsoft YaHei'
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color


def add_text(slide, x, y, w, h, items, fill=None, anchor='top',
             line_spacing=1.15, margin_l=0.04, margin_r=0.04,
             margin_t=0.02, margin_b=0.02):
    if x + w > PAGE_W - 0.05: w = PAGE_W - x - 0.05
    if y + h > PAGE_H - 0.05: h = PAGE_H - y - 0.05
    if x < 0.05: x = 0.05
    if y < 0.05: y = 0.05
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(margin_l)
    tf.margin_right = Inches(margin_r)
    tf.margin_top = Inches(margin_t)
    tf.margin_bottom = Inches(margin_b)
    if anchor == 'middle':
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    elif anchor == 'bottom':
        tf.vertical_anchor = MSO_ANCHOR.BOTTOM
    for i, (t, sz, b, c) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = line_spacing
        p.space_after = Pt(1)
        p.space_before = Pt(1)
        set_run(p.add_run(), t, sz, b, c)
    if fill is not None:
        tb.fill.solid()
        tb.fill.fore_color.rgb = fill
    else:
        tb.fill.background()
    tb.line.fill.background()
    return tb


def add_filled_box(slide, x, y, w, h, fill_color, line_color=None):
    box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(x), Inches(y), Inches(w), Inches(h))
    box.fill.solid()
    box.fill.fore_color.rgb = fill_color
    if line_color is not None:
        box.line.color.rgb = line_color
    else:
        box.line.fill.background()
    return box


def add_oval_marker(slide, x, y, w, h, fill_color, line_color=None):
    sh = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                 Inches(x), Inches(y), Inches(w), Inches(h))
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill_color
    if line_color is not None:
        sh.line.color.rgb = line_color
    else:
        sh.line.fill.background()
    return sh


def add_picture_iso(slide, img_path, x, y, max_w, max_h):
    img = Image.open(img_path)
    iw, ih = img.size
    img_ratio = iw / ih
    calc_h = max_h
    calc_w = calc_h * img_ratio
    if calc_w > max_w:
        calc_w = max_w
        calc_h = calc_w / img_ratio
    actual_x = x + (max_w - calc_w) / 2
    actual_y = y + (max_h - calc_h) / 2
    return slide.shapes.add_picture(img_path,
                                     Inches(actual_x), Inches(actual_y),
                                     width=Inches(calc_w), height=Inches(calc_h))


# ==============================================================
# v3.7.1 修复版
# ==============================================================
# 关键修复：
# 1. 图 1 区域给"正方形"max_h 限制 (2.05 高的盒子里取 1.85 高 + 上下留白)
# 2. 云 BMS 内容放完给 0.02 缓冲，避免和"二、6 项"标题带重叠
# 3. "★核心数据" 单独 textbox 改为 1 个综合 textbox
# 4. 卡片内容上移，"(+1000×)" 行下移到独立位置
# ==============================================================

def build():
    shutil.copy(SRC, OUT)
    prs = Presentation(OUT)
    slide = prs.slides[0]

    for shape in list(slide.shapes):
        sp = shape._element
        sp.getparent().remove(sp)

    # ========== 顶部主标题区 (y=0.05-0.85) ==========
    add_text(slide, 0.20, 0.05, 13.00, 0.50,
             [("0 概念与对比", 14, True, COLOR_HIGHLIGHT),
              ("  ·  ", 14, False, COLOR_LIGHT_GRAY),
              ("云 BMS vs 传统 BMS", 22, True, COLOR_MAIN)],
             line_spacing=1.10)

    add_text(slide, 0.20, 0.55, 13.00, 0.30,
             [("端-边-云协同架构重塑电池管理系统 · ", 11, False, COLOR_GRAY),
              ("算力 1000× · 预警 16 天 · 软件市场 CAGR 20.6%", 11, True, COLOR_RED)],
             line_spacing=1.10)

    # ============================================================
    # 【上排】y=0.95 ~ 3.55
    # 左块：定义 (8.0 宽)       x=0.10 ~ 8.10
    # 右块：图 1 (1:1)          x=10.95 ~ 13.25 (2.30 宽)
    # 右块左：核心数据 (合并 1 个 textbox) x=8.20 ~ 10.90 (2.70 宽)
    # ============================================================

    # ---- 上排左：定义 ----
    add_filled_box(slide, 0.10, 0.95, 8.00, 0.40, COLOR_MAIN)
    add_text(slide, 0.20, 0.95, 7.90, 0.40,
             [("一、定义 · 端-边-云 3 层架构", 13, True, COLOR_WHITE)],
             anchor='middle', line_spacing=1.10)

    # 传统 BMS 小标题
    add_filled_box(slide, 0.10, 1.40, 8.00, 0.30, COLOR_BG_LIGHT)
    add_text(slide, 0.20, 1.40, 7.90, 0.30,
             [("▶ 传统 BMS", 11, True, COLOR_MAIN),
              (" (电池端 MCU)", 9, False, COLOR_LIGHT_GRAY)],
             anchor='middle', line_spacing=1.10)

    add_text(slide, 0.20, 1.72, 7.90, 0.78,
             [("● 单 MCU + AFE · 固定阈值告警 · 无 AI 模型", 10, False, COLOR_GRAY),
              ("● 算力 ~10 MIPS, 存储 KB 级", 10, False, COLOR_GRAY),
              ("● 数据本地化、不出车 · 无远程升级", 10, False, COLOR_GRAY),
              ("● 模式:", 10, True, COLOR_RED),
              (" 被动保护", 10, True, COLOR_RED)])

    # 云 BMS 小标题
    add_filled_box(slide, 0.10, 2.52, 8.00, 0.28, COLOR_BG_RED)
    add_text(slide, 0.20, 2.52, 7.90, 0.28,
             [("▶ 云 BMS", 11, True, COLOR_RED),
              (" (端 + 边 + 云)", 9, False, COLOR_LIGHT_GRAY)],
             anchor='middle', line_spacing=1.10)

    # 云 BMS 内容 (占 y=2.82 ~ 3.55, 留 0.03 缓冲)
    add_text(slide, 0.20, 2.82, 7.90, 0.70,
             [("● 云端算力 + AI 大模型 · 实时状态估计 + 故障预测", 10, False, COLOR_GRAY),
              ("● 算力 ", 10, False, COLOR_GRAY),
              ("10,000+ MIPS", 10, True, COLOR_MAIN),
              (", 存储 TB 级  ", 10, False, COLOR_GRAY),
              ("(+10⁶×)", 10, True, COLOR_RED),
              ("● 数字孪生 + OTA + 跨车队数据闭环 · 主动预测", 10, False, COLOR_GRAY)])

    # ---- 上排右：图 1 雷达图 ----
    add_filled_box(slide, 10.95, 0.95, 2.30, 0.40, COLOR_MAIN)
    add_text(slide, 11.00, 0.95, 2.25, 0.40,
             [("图 1 · 6 维对比", 11, True, COLOR_WHITE)],
             anchor='middle', line_spacing=1.10)

    # 图 1 容器 (y=1.40 ~ 3.55, 高 2.15) - 给图 1.85 高 (留 0.15 上下边距)
    if os.path.exists(CHART_RADAR):
        add_picture_iso(slide, CHART_RADAR,
                        x=10.95, y=1.40, max_w=2.30, max_h=2.05)

    # 核心数据 (合并成 1 个 textbox, y=0.95 ~ 3.55)
    add_filled_box(slide, 8.20, 0.95, 2.70, 2.55, COLOR_BG_BLUE)
    add_text(slide, 8.25, 1.00, 2.60, 2.45,
             [("★ 核心数据", 11, True, COLOR_MAIN),
              ("", 4, False, COLOR_GRAY),
              ("算力提升", 9, False, COLOR_GRAY),
              ("1000×", 14, True, COLOR_RED),
              ("", 3, False, COLOR_GRAY),
              ("存储提升", 9, False, COLOR_GRAY),
              ("10⁶×", 14, True, COLOR_RED),
              ("", 3, False, COLOR_GRAY),
              ("预警提前", 9, False, COLOR_GRAY),
              ("16 天", 14, True, COLOR_RED),
              ("", 3, False, COLOR_GRAY),
              ("市场 CAGR", 9, False, COLOR_GRAY),
              ("20.6%", 14, True, COLOR_RED)],
             line_spacing=1.10, margin_t=0.05, margin_b=0.05)

    # ============================================================
    # 【下排】y=3.65 ~ 7.10
    # ============================================================

    # ---- 下排左：6 项差异 ----
    add_filled_box(slide, 0.10, 3.65, 7.50, 0.40, COLOR_HIGHLIGHT)
    add_text(slide, 0.20, 3.65, 7.40, 0.40,
             [("二、6 项本质差异", 13, True, COLOR_WHITE)],
             anchor='middle', line_spacing=1.10)

    # ★ 修复：卡片高度从 0.95 → 0.85, 留 0.10 垂直间距
    diffs = [
        ("① 算力", "10 MIPS → 10,000+ MIPS", "(+1000×)"),
        ("② 存储", "KB 级 → TB 级", "(+10⁶×)"),
        ("③ 时延", "ms 级本地 → s 级云端协同", ""),
        ("④ 模型更新", "出厂固化 → OTA 持续迭代", ""),
        ("⑤ 业务闭环", "单车数据 → 跨车队", ""),
        ("⑥ 商业模式", "卖硬件 → 软件订阅", ""),
    ]
    card_w = 3.65
    card_h = 0.85
    for i, (label, val, ext) in enumerate(diffs):
        col = i % 2
        row = i // 2
        x = 0.20 + col * (card_w + 0.10)
        y = 4.10 + row * (card_h + 0.10)

        add_filled_box(slide, x, y, card_w, card_h, COLOR_BG_LIGHT)
        add_oval_marker(slide, x + 0.08, y + 0.08, 0.10, 0.10, COLOR_RED)
        # 标题行
        add_text(slide, x + 0.22, y + 0.03, card_w - 0.30, 0.26,
                 [(label, 11, True, COLOR_MAIN)], line_spacing=1.10)
        # 值行
        add_text(slide, x + 0.22, y + 0.30, card_w - 0.30, 0.26,
                 [(val, 9, False, COLOR_GRAY)], line_spacing=1.10)
        # ext 行
        if ext:
            add_text(slide, x + 0.22, y + 0.56, card_w - 0.30, 0.26,
                     [(ext, 9, True, COLOR_RED)], line_spacing=1.10)

    # ---- 下排右：图 2 架构图 ----
    add_filled_box(slide, 7.70, 3.65, 5.55, 0.40, COLOR_HIGHLIGHT)
    add_text(slide, 7.80, 3.65, 5.45, 0.40,
             [("图 2 · 端-边-云 4 能力 × 3 层", 13, True, COLOR_WHITE)],
             anchor='middle', line_spacing=1.10)

    if os.path.exists(CHART_ARCH):
        add_picture_iso(slide, CHART_ARCH,
                        x=7.70, y=4.20, max_w=5.55, max_h=2.85)

    # ========== 底部引言 + 出处 ==========
    add_filled_box(slide, 0.10, 7.20, 13.10, 0.27, COLOR_BG_LIGHT)
    add_text(slide, 0.20, 7.20, 13.00, 0.27,
             [("云 BMS 本质 = ", 8.5, True, COLOR_MAIN),
              ("云端算力 + AI 模型 + 跨资产数据 + 数字孪生", 8.5, True, COLOR_RED),
              ("  |  核心价值:从 ", 8.5, False, COLOR_GRAY),
              ("被动告警", 8.5, True, COLOR_GRAY),
              (" → ", 8.5, False, COLOR_LIGHT_GRAY),
              ("主动预测与价值创造", 8.5, True, COLOR_RED),
              ("  ·  出处:RSC Sustain. Energy Fuels 2025 · 调研报告", 8, False, COLOR_LIGHT_GRAY)],
             anchor='middle', line_spacing=1.10)

    prs.save(OUT)

    # ===== 验证报告 =====
    print("=" * 60)
    print("v3.7.1 修复版 - 验证报告")
    print("=" * 60)

    slide = prs.slides[0]
    print()
    print("【1. 元素可读性 - 图长宽比】")
    pics = [s for s in slide.shapes if s.shape_type == 13]
    for i, shape in enumerate(pics, 1):
        w = shape.width / 914400
        h = shape.height / 914400
        ratio = w/h
        if i == 1:
            status = '✓ 接近 1:1' if 0.85 < ratio < 1.30 else f'✗ {ratio:.2f}'
        else:
            status = '✓ 1.6:1' if 1.4 < ratio < 1.8 else f'✗ {ratio:.2f}'
        print(f'  图 {i}: w={w:.2f} h={h:.2f} ratio={ratio:.2f} {status}')

    print()
    print("【2. 元素可读性 - 文字不堆叠】")
    textboxes = []
    for shape in slide.shapes:
        if shape.has_text_frame and shape.text_frame.text.strip():
            x = shape.left / 914400
            y = shape.top / 914400
            w = shape.width / 914400
            h = shape.height / 914400
            textboxes.append((x, y, x+w, y+h, shape.text_frame.text[:25]))
    overlapping = 0
    for i, (x1, y1, x2, y2, t1) in enumerate(textboxes):
        for x3, y3, x4, y4, t2 in textboxes[i+1:]:
            ox = max(0, min(x2, x4) - max(x1, x3))
            oy = max(0, min(y2, y4) - max(y1, y3))
            if ox * oy > 0.05:
                overlapping += 1
                print(f'  ✗ 重叠: "{t1}" & "{t2}"')
    print(f"  重叠 textbox 数: {overlapping} (应为 0)")

    print()
    print("【3. 空间自适应 - 左/右分配】")
    print("  上排: 左 8.0 (定义) | 右块 8.20-10.90 数据 + 10.95-13.25 图 1")
    print("  下排: 左 7.5 (6 卡) | 右 5.55 (图 2)")

    print()
    print(f"  输出: {OUT}")


if __name__ == "__main__":
    build()
