#!/usr/bin/env python3
"""
extract_template_features.py
============================================================
从多张示例 slide 中归纳"模板共性特征"。

核心思路：
- 3 张 slide 共有的属性 = 模板特征（位置、字号、字体、样式）
- 3 张 slide 各不相同的部分 = 变量（内容、文字）

输出 template_features.json，包含：
- 共有区域（layout_zones）：机构标签区/主标题区/出处区 等
- 共有样式：字号/字体/颜色/加粗
- 变量规则：每个区域下哪些字段是 variable

这是"用户给 3 页示例，技能归纳模板"的核心能力。
"""
import json
import re
import sys
from collections import Counter, defaultdict
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE
from lxml import etree

NS = {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"}


# ============================================================
# 工具函数
# ============================================================

def get_run_color(run):
    """从 run 提取颜色（hex 字符串）"""
    rPr = run._r.find("a:rPr", NS)
    if rPr is not None:
        solidFill = rPr.find("a:solidFill", NS)
        if solidFill is not None:
            srgbClr = solidFill.find("a:srgbClr", NS)
            if srgbClr is not None:
                return srgbClr.get("val")
    return None


def get_run_size(run):
    """从 run 提取字号（pt）"""
    if run.font.size:
        return run.font.size.pt
    rPr = run._r.find("a:rPr", NS)
    if rPr is not None:
        sz = rPr.get("sz")
        if sz:
            return int(sz) / 100
    return None


def get_run_font(run):
    """从 run 提取字体名"""
    if run.font.name:
        return run.font.name
    rPr = run._r.find("a:rPr", NS)
    if rPr is not None:
        latin = rPr.find("a:latin", NS)
        if latin is not None:
            return latin.get("typeface")
    return None


def get_shape_fill_color(shape):
    """从 shape 提取 fill 颜色"""
    try:
        sp = shape._element
        solidFill = sp.find(".//a:solidFill", NS)
        if solidFill is not None:
            srgbClr = solidFill.find("a:srgbClr", NS)
            if srgbClr is not None:
                return srgbClr.get("val")
    except Exception:
        pass
    return None


def get_zone(y_pt, x_pt, slide_w, slide_h):
    """根据坐标判断属于哪个区域
    返回区域 ID（顶部机构区/主标题区/正文区/出处区）
    """
    # 基于相对位置（容差 ±5%）
    rel_y = y_pt / slide_h
    rel_x = x_pt / slide_w

    # 顶部机构区 (y < 0.12)
    if rel_y < 0.12:
        return "header_institution"
    # 主标题区 (0.12 <= y < 0.22) - 放宽到 0.22 以适应 y=54-64
    if rel_y < 0.22:
        return "title_main"
    # 引文/出处区 (y > 0.85)
    if rel_y > 0.85:
        return "footer_source"
    # 左正文区
    if rel_x < 0.55:
        return "body_left"
    # 右图表区
    return "body_right"


def get_role_type(shape):
    """判断 shape 类型"""
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        return "image"
    if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
        return "table"
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        return "group"
    if shape.has_text_frame and shape.text_frame.text.strip():
        return "text"
    return "decoration"


# ============================================================
# 单 slide 特征提取
# ============================================================

def extract_slide_features(slide, slide_w, slide_h):
    """从单张 slide 提取所有形状的特征
    返回 list[dict]：
      {
        zone, role, x, y, w, h,
        runs: [{text, size, bold, color, font}, ...]
      }
    """
    features = []

    def process_shape(shape, zone_override=None):
        left = round(Emu(shape.left).pt, 1) if shape.left else 0
        top = round(Emu(shape.top).pt, 1) if shape.top else 0
        w = round(Emu(shape.width).pt, 1) if shape.width else 0
        h = round(Emu(shape.height).pt, 1) if shape.height else 0

        # 区域判定
        zone = zone_override or get_zone(top, left, slide_w, slide_h)
        role = get_role_type(shape)

        # 文本 run
        runs = []
        if shape.has_text_frame:
            for p_idx, para in enumerate(shape.text_frame.paragraphs):
                for r_idx, run in enumerate(para.runs):
                    runs.append({
                        "p_idx": p_idx,
                        "r_idx": r_idx,
                        "text": run.text,
                        "size": get_run_size(run),
                        "bold": run.font.bold,
                        "color": get_run_color(run),
                        "font": get_run_font(run),
                    })

        # 形状 fill
        fill_color = get_shape_fill_color(shape)

        features.append({
            "zone": zone,
            "role": role,
            "x": left, "y": top, "w": w, "h": h,
            "fill_color": fill_color,
            "runs": runs,
            "full_text": shape.text_frame.text if shape.has_text_frame else "",
        })

        # GROUP 递归
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            for sub in shape.shapes:
                process_shape(sub, zone_override=zone)

    for shape in slide.shapes:
        process_shape(shape)

    return features


# ============================================================
# 多 slide 共性归纳
# ============================================================

def aggregate_features(all_slide_features):
    """从多张 slide 的 features 归纳共性
    思路：
    - 按 zone 分组
    - 统计每个 zone 内最常见的 role、size、bold、color
    - 把"在每张 slide 都有相同 role 的 shape"标记为 template
    """
    # 按 (zone, role) 收集
    zone_role_features = defaultdict(list)

    for slide_features in all_slide_features:
        slide_zones = defaultdict(list)
        for f in slide_features:
            slide_zones[f["zone"]].append(f)
        for zone, fs in slide_zones.items():
            # 每个 zone 多个 shape，按 role 分组
            role_groups = defaultdict(list)
            for f in fs:
                role_groups[f["role"]].append(f)
            for role, group in role_groups.items():
                zone_role_features[(zone, role)].append(group)

    # 找出"每张 slide 都出现"的 (zone, role) 组合
    num_slides = len(all_slide_features)
    common_layouts = []
    for (zone, role), groups in zone_role_features.items():
        if len(groups) == num_slides:  # 每张都有
            # 收集所有 run 特征
            all_runs = []
            for g in groups:
                for f in g:
                    for run in f["runs"]:
                        all_runs.append(run)
            # 统计共性
            common_layouts.append({
                "zone": zone,
                "role": role,
                "count": num_slides,
                "runs": all_runs,
            })

    # === 新增：合并相邻的 zone（title_main 和 body_left 可能因为 y 偏差分裂）===
    # 按特征相似度合并
    merged_zones = _merge_similar_zones(common_layouts, all_slide_features)

    # === 新增：把"主标题特征"的 shape 归到 title_main，无论它在哪个 zone ===
    _reclassify_title_shapes(all_slide_features, merged_zones)

    # 每个 zone 归纳样式
    template_features = {
        "num_slides_analyzed": num_slides,
        "zones": {},
    }

    for layout in merged_zones:
        zone = layout["zone"]
        if zone not in template_features["zones"]:
            template_features["zones"][zone] = {
                "role": layout["role"],
                "occurrences": 0,
                "common_size": None,
                "common_bold": None,
                "common_color": None,
                "common_font": None,
                "x_range": [None, None],
                "y_range": [None, None],
            }

        info = template_features["zones"][zone]
        info["occurrences"] = max(info["occurrences"], layout["count"])

        # 统计 run 级别的共性
        sizes = [r["size"] for r in layout["runs"] if r["size"]]
        bolds = [r["bold"] for r in layout["runs"] if r["bold"] is not None]
        colors = [r["color"] for r in layout["runs"] if r["color"]]
        fonts = [r["font"] for r in layout["runs"] if r["font"]]

        if sizes:
            info["common_size"] = Counter(sizes).most_common(1)[0][0]
        if bolds:
            info["common_bold"] = Counter(bolds).most_common(1)[0][0]
        if colors:
            info["common_color"] = Counter(colors).most_common(1)[0][0]
        if fonts:
            info["common_font"] = Counter(fonts).most_common(1)[0][0]

    # 位置范围
    for slide_features in all_slide_features:
        for f in slide_features:
            zone = f["zone"]
            if zone in template_features["zones"]:
                info = template_features["zones"][zone]
                if info["x_range"][0] is None or f["x"] < info["x_range"][0]:
                    info["x_range"][0] = f["x"]
                if info["x_range"][1] is None or f["x"] > info["x_range"][1]:
                    info["x_range"][1] = f["x"]
                if info["y_range"][0] is None or f["y"] < info["y_range"][0]:
                    info["y_range"][0] = f["y"]
                if info["y_range"][1] is None or f["y"] > info["y_range"][1]:
                    info["y_range"][1] = f["y"]

    # 推断每个 zone 的"关键规则"（哪些是变量）
    _infer_variable_rules(template_features, all_slide_features)

    return template_features


def _merge_similar_zones(common_layouts, all_slide_features):
    """合并相似 zone
    启发：如果两个 zone 的常见样式（size/role）相似，合并
    """
    # 当前：title_main 没出现 3 次（因为 y 偏差），但它的"角色"和样式应该被识别
    # 策略：对出现 ≥2 次的 zone 也加入，作为"近似共性"
    zone_role_features = defaultdict(list)
    for slide_features in all_slide_features:
        slide_zones = defaultdict(list)
        for f in slide_features:
            slide_zones[f["zone"]].append(f)
        for zone, fs in slide_zones.items():
            role_groups = defaultdict(list)
            for f in fs:
                role_groups[f["role"]].append(f)
            for role, group in role_groups.items():
                zone_role_features[(zone, role)].append(group)

    num_slides = len(all_slide_features)
    threshold = max(2, num_slides - 1)  # 至少 2 次或全部 - 1

    merged = []
    for (zone, role), groups in zone_role_features.items():
        if len(groups) >= threshold:
            all_runs = []
            for g in groups:
                for f in g:
                    for run in f["runs"]:
                        all_runs.append(run)
            merged.append({
                "zone": zone,
                "role": role,
                "count": len(groups),
                "runs": all_runs,
            })

    return merged


def _reclassify_title_shapes(all_slide_features, merged_zones):
    """重分类：把"主标题特征"的 shape 归到 title_main
    主标题特征：22pt 加粗 或 ≥ 20pt
    """
    for slide_features in all_slide_features:
        for f in slide_features:
            if not f["runs"]:
                continue
            # 检测是否是主标题
            is_title = False
            for run in f["runs"]:
                if run["size"] and run["size"] >= 20 and run["bold"]:
                    is_title = True
                    break
            if is_title and f["zone"] != "title_main":
                # 重新归类
                old_zone = f["zone"]
                f["zone"] = "title_main"
                # 同步更新 merged_zones 中的统计
                for layout in merged_zones:
                    if layout["zone"] == old_zone:
                        # 从原 zone 移除
                        layout["runs"] = [r for r in layout["runs"] if r not in f["runs"]]
                    if layout["zone"] == "title_main":
                        # 添加到 title_main
                        for run in f["runs"]:
                            if run not in layout["runs"]:
                                layout["runs"].append(run)
                                layout["count"] += 1


def _infer_variable_rules(template_features, all_slide_features):
    """推断每个 zone 下的变量规则
    启发式：
    - 文本内容是变量（每张 slide 不同）
    - 字号/字体/颜色/位置是模板（在每张 slide 一致）
    - 加粗且颜色为红色的 run 标记为"关键强调"格式
    """
    for zone, info in template_features["zones"].items():
        # 收集所有 slide 中该 zone 的文本
        all_texts = []
        for slide_features in all_slide_features:
            for f in slide_features:
                if f["zone"] == zone and f["full_text"]:
                    all_texts.append(f["full_text"])

        info["variable_rules"] = {
            "text_is_variable": len(set(all_texts)) > 1,  # 文本不同 = 变量
            "format_is_template": True,  # 字号/颜色/位置是模板
            "emphasis_pattern": _detect_emphasis_pattern(all_texts, info),
        }


def _detect_emphasis_pattern(texts, info):
    """检测"红色加粗强调"模式
    例如：标题中部分 run 是红色加粗 → 强调
    """
    common_color = info.get("common_color")
    common_size = info.get("common_size") or 0
    if common_color in ("FF0000", "C00000"):  # 红色
        return "title_with_red_emphasis"
    if common_size >= 20:
        return "main_title"
    if common_size >= 10 and common_size <= 14:
        return "body_text"
    return "generic"


# ============================================================
# 主入口
# ============================================================

def extract_template_features(pptx_path):
    """主入口"""
    prs = Presentation(pptx_path)
    slide_w = prs.slide_width.pt
    slide_h = prs.slide_height.pt

    all_slide_features = []
    for slide in prs.slides:
        features = extract_slide_features(slide, slide_w, slide_h)
        all_slide_features.append(features)

    template_features = aggregate_features(all_slide_features)
    template_features["source"] = str(Path(pptx_path).resolve())
    template_features["slide_size"] = {"width": slide_w, "height": slide_h}

    return template_features


def main():
    if len(sys.argv) != 3:
        print("Usage: python extract_template_features.py <template.pptx> <output.json>")
        sys.exit(1)

    pptx_path = sys.argv[1]
    output_path = sys.argv[2]

    features = extract_template_features(pptx_path)
    with open(output_path, "w", encoding="utf-8") as f:
        json.dump(features, f, ensure_ascii=False, indent=2)

    print(f"✅ 模板特征提取完成：{output_path}")
    print(f"   分析的 slide 数：{features['num_slides_analyzed']}")
    print(f"   识别到的共性 zone 数：{len(features['zones'])}")
    for zone, info in features["zones"].items():
        print(f"     {zone}: role={info['role']}, size={info['common_size']}, "
              f"color={info['common_color']}, occurrences={info['occurrences']}")


if __name__ == "__main__":
    main()
