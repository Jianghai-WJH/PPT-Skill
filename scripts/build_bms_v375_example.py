"""
v3.7.4 修复：
1. 上排内容溢出到下排 (云 BMS 内容 h=0.70 → 0.55)
2. 卡片内值/ext 重叠 0.01"
3. 所有 textbox 严格不重叠
"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from PIL import Image

PAGE_W, PAGE_H = 13.33, 7.50

COLOR_TITLE_BG = RGBColor(0x0E, 0x3A, 0x6B)
COLOR_RED = RGBColor(0xC0, 0x39, 0x2B)
COLOR_BAND_BLUE = RGBColor(0x1F, 0x4E, 0x79)
COLOR_BAND_ORANGE = RGBColor(0xCC, 0x6E, 0x1F)
COLOR_CARD_BG = RGBColor(0xF5, 0xF7, 0xFA)
COLOR_TEXT = RGBColor(0x2C, 0x3E, 0x50)
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)


def add_text(slide, x, y, w, h, text, font_pt=10, bold=False, color=COLOR_TEXT,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, bg_color=None, line_spacing=1.20):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.03); tf.margin_right = Inches(0.03)
    tf.margin_top = Inches(0.02); tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = anchor
    if bg_color is not None:
        tb.fill.solid(); tb.fill.fore_color.rgb = bg_color
        tb.line.fill.background()
    lines = text.split('\n')
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align; p.line_spacing = line_spacing
        run = p.add_run(); run.text = line
        run.font.size = Pt(font_pt); run.font.bold = bold
        run.font.color.rgb = color; run.font.name = 'Microsoft YaHei'
    return tb


def add_band(slide, x, y, w, h, text, color, font_pt=13, font_color=COLOR_WHITE):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid(); shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    tf = shape.text_frame
    tf.margin_left = Inches(0.1); tf.margin_right = Inches(0.05)
    tf.margin_top = Inches(0.04); tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    run = p.add_run(); run.text = text
    run.font.size = Pt(font_pt); run.font.bold = True
    run.font.color.rgb = font_color; run.font.name = 'Microsoft YaHei'


def add_picture_iso(slide, img_path, x, y, max_w, max_h):
    img = Image.open(img_path)
    iw, ih = img.size
    img_ratio = iw / ih
    calc_h = max_h
    calc_w = calc_h * img_ratio
    if calc_w > max_w:
        calc_w = max_w
        calc_h = calc_w / img_ratio
    actual_x = x + (max_w - calc_w) / 2
    actual_y = y + (max_h - calc_h) / 2
    return slide.shapes.add_picture(img_path, Inches(actual_x), Inches(actual_y),
                                     Inches(calc_w), Inches(calc_h))


# ============= 主程序 =============
prs = Presentation()
prs.slide_width = Inches(PAGE_W); prs.slide_height = Inches(PAGE_H)
blank = prs.slide_layouts[6]
slide = prs.slides.add_slide(blank)

bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(PAGE_W), Inches(PAGE_H))
bg.fill.solid(); bg.fill.fore_color.rgb = COLOR_WHITE
bg.line.fill.background()

# 1. 主标题 (h=1.20)
TITLE_H = 1.20
add_text(slide, 0.20, 0.05, 13.00, TITLE_H - 0.05, "0  概念与对比\n云 BMS vs 传统 BMS",
         font_pt=22, bold=True, color=COLOR_TITLE_BG, line_spacing=1.10)

# 副标题红色横条
sub_y = 0.05 + TITLE_H + 0.05
SUB_H = 0.40
sub_band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.20), Inches(sub_y),
                                    Inches(13.00), Inches(SUB_H))
sub_band.fill.solid(); sub_band.fill.fore_color.rgb = COLOR_RED
sub_band.line.fill.background()
tf = sub_band.text_frame
tf.margin_left = Inches(0.15); tf.margin_top = Inches(0.05)
tf.vertical_anchor = MSO_ANCHOR.MIDDLE
p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
run = p.add_run()
run.text = "端-边-云协同架构重塑电池管理系统  ·  算力 1000×  ·  预警 16 天  ·  软件市场 CAGR 20.6%"
run.font.size = Pt(11); run.font.bold = True
run.font.color.rgb = COLOR_WHITE; run.font.name = 'Microsoft YaHei'

# 2. 上排左：定义 (压缩到 2.35 高)
UPPER_Y = sub_y + SUB_H + 0.20  # 1.90
UPPER_H = 2.40
LEFT_W = 7.50
LEFT_X = 0.20

add_band(slide, LEFT_X, UPPER_Y, LEFT_W, 0.45,
         "一、定义 · 端-边-云 3 层架构", COLOR_BAND_BLUE, font_pt=14)

# 传统 BMS: 子标题 0.30 + 内容 0.55 (3行 × 10pt × 1.20 = 0.50 + padding)
trad_y = UPPER_Y + 0.50
add_text(slide, LEFT_X + 0.05, trad_y, LEFT_W - 0.10, 0.30, "▶ 传统 BMS（电池端 MCU）",
         font_pt=12, bold=True, color=COLOR_BAND_BLUE, line_spacing=1.20)
trad_content_y = trad_y + 0.34
add_text(slide, LEFT_X + 0.05, trad_content_y, LEFT_W - 0.10, 0.60,
         "● 单 MCU + AFE · 固定阈值告警\n● 算力 ~10 MIPS, 存储 KB 级\n● 无 AI 模型，无云端协同",
         font_pt=10, color=COLOR_TEXT, line_spacing=1.20)

# 云 BMS: 子标题 0.30 + 内容 0.55
cloud_y = trad_content_y + 0.55 + 0.05  # 2.27 + 0.32 + 0.55 + 0.05 = 2.95... 重新算
# trad_y = UPPER_Y + 0.50 = 1.90+0.50 = 2.40
# trad_content_y = 2.40 + 0.32 = 2.72
# cloud_y = 2.72 + 0.55 + 0.05 = 3.32
# cloud_content_y = 3.32 + 0.32 = 3.64
# cloud_content 结束 = 3.64 + 0.55 = 4.19
# UPPER_Y + UPPER_H = 1.90 + 2.35 = 4.25 ✓ 不超
add_text(slide, LEFT_X + 0.05, cloud_y, LEFT_W - 0.10, 0.30, "▶ 云 BMS（端 + 边 + 云）",
         font_pt=12, bold=True, color=COLOR_RED, line_spacing=1.20)
add_text(slide, LEFT_X + 0.05, cloud_y + 0.32, LEFT_W - 0.10, 0.60,
         "● 云端算力 + AI 大模型 · 实时状态估计 + 故障预测\n● 算力 10,000+ MIPS, 存储 TB 级\n● 跨车队数据闭环，OTA 持续迭代",
         font_pt=10, color=COLOR_TEXT, line_spacing=1.20)

# 3. 上排右：核心数据
RIGHT_X = 7.80
RIGHT_W = 5.40
add_text(slide, RIGHT_X, UPPER_Y, RIGHT_W, 0.30, "★ 核心数据",
         font_pt=14, bold=True, color=COLOR_RED, line_spacing=1.20)
data_items = [
    ("算力提升", "1000×", "10 → 10,000+ MIPS"),
    ("存储提升", "10⁶×", "KB → TB 级"),
    ("预警提前", "16 天", "热失控预测"),
    ("市场 CAGR", "20.6%", "2026-2032"),
]
item_h = 0.42
for i, (label, big, sub) in enumerate(data_items):
    iy = UPPER_Y + 0.35 + i * item_h
    add_text(slide, RIGHT_X, iy, 1.5, item_h, label,
             font_pt=10, bold=True, color=COLOR_BAND_BLUE, line_spacing=1.10)
    add_text(slide, RIGHT_X + 1.5, iy, 1.5, item_h, big,
             font_pt=15, bold=True, color=COLOR_RED, line_spacing=1.10, align=PP_ALIGN.CENTER)
    add_text(slide, RIGHT_X + 3.0, iy, 2.4, item_h, sub,
             font_pt=9, color=COLOR_TEXT, line_spacing=1.10)

# 4. 下排左：6 项差异 (修复卡片)
LOWER_Y = UPPER_Y + UPPER_H + 0.15  # 4.40
LOWER_H = 2.35
LOWER_LEFT_W = 7.50
add_band(slide, LEFT_X, LOWER_Y, LOWER_LEFT_W, 0.40,
         "二、6 项本质差异", COLOR_BAND_ORANGE, font_pt=13)

cards_y = LOWER_Y + 0.45
card_h = 0.62
card_w = 3.65
card_gap_x = 0.10
card_gap_y = 0.05

cards = [
    ("① 算力", "10 MIPS → 10,000+ MIPS", "(+1000×)"),
    ("② 存储", "KB 级 → TB 级", "(+10⁶×)"),
    ("③ 时延", "ms 级本地 → s 级云端协同", ""),
    ("④ 模型更新", "出厂固化 → OTA 持续迭代", ""),
    ("⑤ 业务闭环", "单车数据 → 跨车队", ""),
    ("⑥ 商业模式", "卖硬件 → 软件订阅", ""),
]

for idx, (label, val, ext) in enumerate(cards):
    row = idx // 2
    col = idx % 2
    cx = LEFT_X + col * (card_w + card_gap_x)
    cy = cards_y + row * (card_h + card_gap_y)
    card_bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(cx), Inches(cy),
                                      Inches(card_w), Inches(card_h))
    card_bg.fill.solid(); card_bg.fill.fore_color.rgb = COLOR_CARD_BG
    card_bg.line.color.rgb = COLOR_BAND_ORANGE
    card_bg.line.width = Pt(0.5)
    # 修复：标签 0.22 / 值 0.18 (y=0.25) / ext 0.16 (y=0.44, 间距 0.01→0.05)
    add_text(slide, cx + 0.05, cy + 0.03, card_w - 0.10, 0.22, label,
             font_pt=11, bold=True, color=COLOR_BAND_ORANGE, line_spacing=1.10)
    add_text(slide, cx + 0.05, cy + 0.26, card_w - 0.10, 0.18, val,
             font_pt=9, color=COLOR_TEXT, line_spacing=1.10)
    if ext:
        add_text(slide, cx + 0.05, cy + 0.45, card_w - 0.10, 0.16, ext,
                 font_pt=8, bold=True, color=COLOR_RED, line_spacing=1.10)

# 5. 下排右：雷达图 + 架构图
LOWER_RIGHT_X = 7.80
LOWER_RIGHT_W = 5.40
add_text(slide, LOWER_RIGHT_X, LOWER_Y, LOWER_RIGHT_W, 0.25, "图 1 · 6 维评分对比",
         font_pt=11, bold=True, color=COLOR_BAND_BLUE, align=PP_ALIGN.CENTER, line_spacing=1.10)
add_picture_iso(slide, 'artifacts/charts/01_雷达图_6维对比.png',
                LOWER_RIGHT_X, LOWER_Y + 0.27, LOWER_RIGHT_W, 1.05)

arch_y = LOWER_Y + 1.40
add_text(slide, LOWER_RIGHT_X, arch_y, LOWER_RIGHT_W, 0.25, "图 2 · 端-边-云 4 能力 × 3 层",
         font_pt=11, bold=True, color=COLOR_BAND_ORANGE, align=PP_ALIGN.CENTER, line_spacing=1.10)
add_picture_iso(slide, 'artifacts/charts/06_架构图_端边云.png',
                LOWER_RIGHT_X, arch_y + 0.27, LOWER_RIGHT_W, 1.10)

# 6. 底部引言
BOTTOM_Y = LOWER_Y + LOWER_H + 0.10  # 4.40+2.40+0.10 = 6.90
BOTTOM_H = 0.40
add_text(slide, 0.20, BOTTOM_Y, 13.00, BOTTOM_H,
         "云 BMS 本质 = 云端算力 + AI 模型 + 跨资产数据 + 数字孪生  ·  出处: RSC 2025 / 调研报告",
         font_pt=10, color=COLOR_TEXT, line_spacing=1.20, bg_color=COLOR_CARD_BG)

out = 'artifacts/BMS_vs_云BMS_对比页_v3.7.5_零堆叠最终版.pptx'
prs.save(out)
print(f'OK Saved: {out}')
print(f'  Total slide height used: {BOTTOM_Y + BOTTOM_H:.2f} / {PAGE_H}')
