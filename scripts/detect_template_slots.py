#!/usr/bin/env python3
"""
detect_template_slots.py
============================================================
核心组件：把 PPT 中每个 run/段落自动分类为 template(模板保留)
或 variable(变量待替换)。

支持两种模式：
  1) 启发式(heuristic) ：纯规则，适合大多数场景
  2) AI 提示(ai-hint)  ：生成给 AI 看的判定清单，让 AI 决策

判定粒度：
  - run 级    ：每个 <a:r> 单独判定
  - 段落级    ：空段落、装饰性段落
  - 形状级    ：图片/装饰/占位符
  - 整 slide  : 决定保留哪些形状、删除哪些形状

输出 slot_map.json，内容生成阶段直接消费。
"""
import json
import re
import sys
from collections import Counter
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE


# ============================================================
# 启发式判定规则
# ============================================================

# 常见"模板标签"模式：序号、章节号、单位符号、连接词
LABEL_PATTERNS = [
    r"^\d+(\.\d+)*$",  # 纯数字 / 多级编号 "3" "3.3"
    r"^\d+、$",  # 中文编号 "1、"
    r"^（\d+）$",  # 中文括号编号 "（1）"
    r"^\(\d+\)$",  # 英文括号编号
    r"^[一二三四五六七八九十]+、$",  # 中文数字编号
    r"^[A-Z]+\d+$",  # 字母+数字
    r"^#",  # markdown 标题符号遗留
]

# 常见"模板小标签"（5W2H、问题分类标签等）
FIXED_LABELS = {
    "5W2H", "5w2h",
    "When", "Where", "Who", "What", "Why", "How", "How many",
    "SWOT", "PEST", "PDCA",
}

# 常见"模板小标题词"（在公司汇报 PPT 中通常作为可保留的小标题）
TEMPLATE_HEADING_HINTS = [
    "故障现象", "原因分析", "原因解析", "对策", "临时对策", "长久对策",
    "对策实施计划", "实施计划", "总结", "结论", "问题描述", "背景",
    "解决方案", "下一步", "风险", "结论与展望", "回顾", "目标",
    "5W2H", "故障复盘", "问题复盘",
]

# 常见"装饰性 / 非内容"占位符
DECORATION_KEYWORDS = [
    "logo", "Logo", "LOGO",
    "标识", "图标",
    "装饰", "background", "Background",
    "页码", "页眉", "页脚", "footer", "header",
    "水印", "watermark",
]

# 数字/单位/符号串（更可能是数据点而非模板）
# 注意：纯数字（如 "20"、"5"）也属于数据，应该标记为 variable
DATA_LIKE_PATTERN = re.compile(
    r"^[\d\.\,\-\+\s℃°%‰×÷±≤≥<>=/\*]+(kw|KW|kW|V|A|Hz|rpm|MPa|kPa|N·m|°C|℃|mm|cm|m|km/h|L|min|s|h|mg|g|kg|t|RMB|元|万|亿|%)?$"
)
# 纯数字串：日期、月份、分钟数、计数等
PURE_NUMBER_PATTERN = re.compile(r"^[\d\.\,]+$")


def is_label_like(text):
    """判断一个 run 文本是否更像'模板标签'（编号、序号、章节号）"""
    text = text.strip()
    if not text:
        return False
    for pat in LABEL_PATTERNS:
        if re.match(pat, text):
            return True
    return False


def is_data_like(text):
    """判断一个 run 文本是否像数据点（纯数字、单位串）"""
    text = text.strip()
    if not text:
        return False
    if DATA_LIKE_PATTERN.match(text):
        return True
    if PURE_NUMBER_PATTERN.match(text):
        # 纯数字串（"20"、"5"、"3.3" 单看像章节号，但若是孤立的也是数据）
        # 但要排除像 "3" 这样的简单编号（已经被 is_label_like 捕获）
        return True
    return False


def is_chapter_like(text):
    """判断是否是"章节号"模板（"3"、"3.3"、"3.3.1"）
    与"纯数字数据"区分：带小数点的多级编号算章节号
    """
    text = text.strip()
    if not text:
        return False
    # 多级编号才视为章节号
    if re.match(r"^\d+(\.\d+)+$", text):
        return True
    return False


def is_punctuation_only(text):
    text = text.strip()
    if not text:
        return False
    return all(c in "：:、,，.。;；!?！？()【】[]「」"" '' +-*/=×÷—…《》<>·•" for c in text)


def is_template_label_word(text):
    """判断是否是固定模板小标签词"""
    text = text.strip()
    return text in FIXED_LABELS


def looks_like_template_heading(text, all_texts):
    """判断 run 文本是否是模板小标题
    启发：若一段(后续无冒号紧跟内容)的文字+后面跟着若干普通段落，
    且该文字本身短小、加粗、字号偏大 → 模板小标题
    """
    text = text.strip()
    if not text or len(text) > 30:
        return False
    # 命中模板小标题词表
    for hint in TEMPLATE_HEADING_HINTS:
        if hint in text or text in hint:
            return True
    # 末尾是冒号的小标题（"5W2H:" "原因分析："）
    if text.endswith("：") or text.endswith(":"):
        stem = text.rstrip("：:").strip()
        if stem and len(stem) <= 12:
            # 短小带冒号，更像模板
            return True
    return False


# ============================================================
# 形状分类
# ============================================================

def classify_shape_role(shape):
    """分类一个形状的角色
    返回: 'placeholder_title' | 'placeholder_body'
         | 'decoration' | 'logo' | 'image'
         | 'text_label' | 'data_block' | 'unknown'
    """
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        return "image"
    if shape.is_placeholder:
        ph_type = str(shape.placeholder_format.type) if shape.placeholder_format.type else ""
        if "TITLE" in ph_type:
            return "placeholder_title"
        if "BODY" in ph_type or "CONTENT" in ph_type:
            return "placeholder_body"
        return "placeholder_other"
    name_l = (shape.name or "").lower()
    for kw in DECORATION_KEYWORDS:
        if kw.lower() in name_l:
            return "decoration"
    if shape.has_text_frame and shape.text_frame.text.strip():
        return "text_label"
    return "decoration"


# ============================================================
# 逐 run 分类
# ============================================================

def classify_run_in_paragraph(p_idx, para, paragraph_context):
    """对一个段落中的每个 run 分类为 'template' 或 'variable'

    paragraph_context 提供该段落的前后文（用于启发式判断）
    """
    classifications = []
    for r_idx, run in enumerate(para.runs):
        text = run.text
        stripped = text.strip()
        font = run.font

        # 决定 run 角色
        role = "variable"  # 默认是变量

        if not stripped:
            # 空 run：保留为 template（一般是装饰/格式需要）
            role = "template"
        elif is_punctuation_only(stripped):
            # 纯标点：模板
            role = "template"
        elif is_chapter_like(stripped):
            # 多级章节号（"3.3"、"3.3.1"）：模板
            role = "template"
        elif is_label_like(stripped) and not is_data_like(stripped):
            # 编号/序号：模板
            role = "template"
        elif is_data_like(stripped):
            # 纯数据/单位/数字：变量（属于用户要替换的内容）
            role = "variable"
        elif is_template_label_word(stripped):
            # 固定模板标签词（5W2H/When 等）
            role = "template"
        elif looks_like_template_heading(stripped, paragraph_context):
            # 模板小标题
            role = "template"
        else:
            # 其他普通文本：变量
            role = "variable"

        classifications.append({
            "p_idx": p_idx,
            "r_idx": r_idx,
            "text": text,
            "role": role,
            "font_size_pt": float(font.size.pt) if font.size else None,
            "bold": font.bold,
            "font_name": font.name,
        })
    return classifications


def detect_slots_in_shape(shape, shape_idx, shape_role):
    """检测一个形状内的所有 run 角色
    返回 slot 列表（每个 slot 对应一个 run 或一个段落）

    特殊处理：
    - GROUP：递归处理内部 sub-shapes（每组一个 slots 列表）
    - PICTURE：无 slot（保留图片或删除）
    - 普通 TEXT：逐 run 分类
    """
    # 组合：递归处理内部
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        group_slots = []
        for sub_idx, sub_shape in enumerate(shape.shapes):
            sub_role = classify_shape_role(sub_shape)
            sub_slots = detect_slots_in_shape(sub_shape, sub_idx, sub_role)
            group_slots.append({
                "sub_idx": sub_idx,
                "sub_name": sub_shape.name,
                "sub_role": sub_role,
                "slots": sub_slots,
            })
        return [{
            "shape_idx": shape_idx,
            "shape_name": shape.name,
            "shape_role": shape_role,
            "kind": "group",
            "sub_shapes": group_slots,
            "default_action": "edit_runs",  # 组合内含可编辑文字，标记为 edit
        }]

    # 图片：无 slot
    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
        return [{
            "shape_idx": shape_idx,
            "shape_name": shape.name,
            "shape_role": shape_role,
            "kind": "image",
            "default_role": "template",  # 保留 logo 图片由用户控制
        }]

    # 无文本框
    if not shape.has_text_frame:
        return [{
            "shape_idx": shape_idx,
            "shape_name": shape.name,
            "shape_role": shape_role,
            "kind": "shape_keep",
            "default_role": "template",
            "runs": [],
            "full_text": "",
        }]

    tf = shape.text_frame
    # 收集整 shape 文本（用于上下文）
    all_texts = []
    for p in tf.paragraphs:
        for r in p.runs:
            all_texts.append(r.text)
    paragraph_context = " ".join(all_texts)

    slots = []
    for p_idx, para in enumerate(tf.paragraphs):
        run_classifications = classify_run_in_paragraph(p_idx, para, paragraph_context)
        slots.append({
            "shape_idx": shape_idx,
            "shape_name": shape.name,
            "shape_role": shape_role,
            "kind": "paragraph",
            "p_idx": p_idx,
            "runs": run_classifications,
            "full_text": "".join(r.text for r in para.runs),
        })

    return slots


# ============================================================
# 主流程
# ============================================================

def detect_slots(pptx_path, mode="heuristic"):
    """检测模板中的 template/variable slots

    mode:
      - "heuristic"   : 纯规则判定
      - "ai-hint"     : 输出给 AI 看的清单，让 AI 决策
    """
    prs = Presentation(pptx_path)

    # 收集所有 slide 的形状和 slot
    slides_info = []
    for s_idx, slide in enumerate(prs.slides):
        slide_info = {
            "slide_idx": s_idx,
            "layout": slide.slide_layout.name,
            "shapes": [],
        }

        for sh_idx, shape in enumerate(slide.shapes):
            role = classify_shape_role(shape)
            slots = detect_slots_in_shape(shape, sh_idx, role)

            shape_info = {
                "sh_idx": sh_idx,
                "name": shape.name,
                "type": str(shape.shape_type) if shape.shape_type else None,
                "role": role,
                "position": {
                    "left_pt": round(Emu(shape.left).pt, 1) if shape.left else None,
                    "top_pt": round(Emu(shape.top).pt, 1) if shape.top else None,
                    "width_pt": round(Emu(shape.width).pt, 1) if shape.width else None,
                    "height_pt": round(Emu(shape.height).pt, 1) if shape.height else None,
                },
                "slots": slots,
                "default_action": _default_action_for_role(role),
            }
            slide_info["shapes"].append(shape_info)

        slides_info.append(slide_info)

    # 汇总：每个 shape 的最终决策
    decision = _build_decision(slides_info, mode)

    return {
        "source": str(Path(pptx_path).resolve()),
        "mode": mode,
        "slides": slides_info,
        "decision": decision,
    }


def _default_action_for_role(role):
    """基于形状角色的默认动作"""
    if role == "image":
        return "delete"  # 模板里的图片默认删除（用户可重新插入）
    if role == "decoration":
        return "delete"  # 装饰默认删除
    if role == "logo":
        return "keep"  # logo 保留
    if role == "placeholder_title":
        return "edit_runs"  # 标题占位符：按 run 编辑
    if role == "placeholder_body":
        return "edit_runs"  # 正文占位符：按 run 编辑
    if role == "text_label":
        return "edit_runs"
    return "edit_runs"


def _build_decision(slides_info, mode):
    """生成 AI 友好的判定清单（如果 mode=ai-hint）
    或最终决策（如果 mode=heuristic）
    """
    decisions = []
    for slide in slides_info:
        slide_decision = {
            "slide_idx": slide["slide_idx"],
            "layout": slide["layout"],
            "shapes": [],
        }

        for shape in slide["shapes"]:
            shape_decision = {
                "sh_idx": shape["sh_idx"],
                "name": shape["name"],
                "role": shape["role"],
                "default_action": shape["default_action"],
                "slot_summary": _summarize_slots(shape["slots"]),
            }

            if mode == "ai-hint":
                # 给 AI 看的：列出每个 run + 启发式判定 + 候选解释
                shape_decision["hint_for_ai"] = _build_ai_hint(shape)
            else:
                # heuristic 模式：直接给最终判定
                shape_decision["final_runs"] = _finalize_runs(shape["slots"])

            slide_decision["shapes"].append(shape_decision)
        decisions.append(slide_decision)
    return decisions


def _summarize_slots(slots):
    """汇总一个 shape 的 slot 统计"""
    if not slots:
        return {"template_count": 0, "variable_count": 0, "kind": "empty"}

    template_count = 0
    variable_count = 0
    kind = "shape_only"

    def count_runs(slot_list):
        nonlocal template_count, variable_count
        for slot in slot_list:
            if slot["kind"] == "paragraph":
                for run in slot["runs"]:
                    if run["role"] == "template":
                        template_count += 1
                    else:
                        variable_count += 1
            elif slot["kind"] == "group":
                # 递归统计组合内
                for sub in slot["sub_shapes"]:
                    count_runs(sub["slots"])

    for slot in slots:
        if slot["kind"] == "paragraph":
            kind = "paragraph"
        elif slot["kind"] == "group":
            kind = "group"

    count_runs(slots)
    return {
        "template_count": template_count,
        "variable_count": variable_count,
        "kind": kind,
    }


def _build_ai_hint(shape):
    """为 AI 构造判定提示
    AI 拿到这个提示后，可以逐 run 确认/修改 role
    """
    hint_lines = []
    hint_lines.append(f"Shape: {shape['name']} (role={shape['role']})")
    hint_lines.append("请逐 run 判定每个 run 是 'template'(保留) 还是 'variable'(可替换):")

    def walk_slots(slot_list, prefix=""):
        for slot in slot_list:
            if slot["kind"] == "paragraph":
                for run in slot["runs"]:
                    hint_lines.append(
                        f"  {prefix}P{slot['p_idx']}R{run['r_idx']} role={run['role']:9s} | "
                        f"text=\"{run['text'][:40]}\" "
                        f"size={run['font_size_pt']} bold={run['bold']}"
                    )
            elif slot["kind"] == "group":
                for sub in slot["sub_shapes"]:
                    hint_lines.append(f"  {prefix}[Group: {sub['sub_name']}]")
                    walk_slots(sub["slots"], prefix + "  ")

    walk_slots(shape["slots"])
    hint_lines.append(
        "判定规则: 章节号/编号/标点/小标题词 → template; 具体描述/数据/单位 → variable"
    )
    return "\n".join(hint_lines)


def _finalize_runs(slots):
    """heuristic 模式下：把每个 run 的 role 直接作为最终判定
    重要：path 中包含 sub_shape 路径（如果在组合内）
    """
    final = []

    def walk(slot_list, path_prefix=""):
        for slot in slot_list:
            if slot["kind"] == "paragraph":
                for run in slot["runs"]:
                    path = f"{path_prefix}p{run['p_idx']}.r{run['r_idx']}"
                    final.append({
                        "p_idx": run["p_idx"],
                        "r_idx": run["r_idx"],
                        "path": path,
                        "text": run["text"],
                        "role": run["role"],
                    })
            elif slot["kind"] == "group":
                for sub in slot["sub_shapes"]:
                    sub_prefix = f"{path_prefix}sub{sub['sub_idx']}."
                    walk(sub["slots"], sub_prefix)

    walk(slots)
    return final


def main():
    if len(sys.argv) < 3:
        print("Usage: python detect_template_slots.py <template.pptx> <output.json> [--mode heuristic|ai-hint]")
        sys.exit(1)

    pptx_path = sys.argv[1]
    output_path = sys.argv[2]

    mode = "heuristic"
    if "--mode" in sys.argv:
        idx = sys.argv.index("--mode")
        mode = sys.argv[idx + 1]

    result = detect_slots(pptx_path, mode=mode)
    with open(output_path, "w", encoding="utf-8") as f:
        json.dump(result, f, ensure_ascii=False, indent=2)

    print(f"✅ Slot 检测完成：{output_path} (mode={mode})")
    # 简明汇总
    total_template = 0
    total_variable = 0
    for slide in result["decision"]:
        for shape in slide["shapes"]:
            s = shape["slot_summary"]
            total_template += s["template_count"]
    